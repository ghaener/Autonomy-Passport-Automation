from pptx import Presentation
from pptx.util import Inches, Pt

# Ouvre le fichier contenant les prénom/noms
with open('Liste nom.txt', 'r') as file:
    text = file.read().splitlines()    
    for lines in text:
        # Set variable nom & prenom
        # A ADAPTER EN FONCTION DU FORMAT DE LA LISTE
        nom = lines.split()[0]
        prenom = lines.split()[1]
        
        #Ouvre le template
        prs = Presentation('GRIMPO6-PasseportAutonomie_2025_Adulte.pptx')
        for slide in prs.slides: 
           
            # Add textbox nom
            txt_box = slide.shapes.add_textbox(Inches(1.85), Inches(0.48), Inches(3), Inches(3))
            txt_frame = txt_box.text_frame
            txt_frame.text = nom
            
            # Add textbox prenom
            txt_box = slide.shapes.add_textbox(Inches(2.1), Inches(0.85), Inches(3), Inches(3))
            txt_frame = txt_box.text_frame
            txt_frame.text = prenom

            #Sauvegarde un nouveau fichier
            prs.save('Passport Liste/Passport_' + nom + '_' + prenom + '.pptx')